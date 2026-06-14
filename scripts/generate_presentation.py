import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("[ERROR] python-pptx is not installed. Run: pip install python-pptx")
    sys.exit(1)

def create_presentation():
    prs = Presentation()
    # Set standard 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette Constants
    DARK_BG = RGBColor(15, 23, 42)      # Slate 900 (Deep rich dark)
    LIGHT_TEXT = RGBColor(241, 245, 249) # Slate 100
    MUTED_TEXT = RGBColor(148, 163, 184) # Slate 400
    ACCENT_GREEN = RGBColor(16, 185, 129) # Emerald 500 (Authorized/Success)
    ACCENT_BLUE = RGBColor(59, 130, 246)  # Blue 500 (Processing/AI)
    ACCENT_RED = RGBColor(239, 68, 68)    # Red 500 (Unknown/Alert)

    # Blank layout for custom styling
    blank_layout = prs.slide_layouts[6]

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BG

    def add_title(slide, text, color=LIGHT_TEXT):
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Calibri'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color
        return txBox

    # ==========================================
    # SLIDE 1: Title Slide (Dark Theme)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1)
    
    # Title Box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "NVIDIA DEEPSTREAM SURVEILLANCE PIPELINE"
    p.font.name = 'Calibri'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    
    p2 = tf.add_paragraph()
    p2.text = "Restructured Isolated Python & C++ Architectures"
    p2.font.name = 'Calibri'
    p2.font.size = Pt(26)
    p2.font.color.rgb = LIGHT_TEXT
    
    # Subtitle
    sub_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.5))
    tf_sub = sub_box.text_frame
    p_sub1 = tf_sub.paragraphs[0]
    p_sub1.text = "Ha-Meem Group AI Surveillance Engineering"
    p_sub1.font.name = 'Calibri'
    p_sub1.font.size = Pt(18)
    p_sub1.font.bold = True
    p_sub1.font.color.rgb = ACCENT_BLUE
    
    p_sub2 = tf_sub.add_paragraph()
    p_sub2.text = "Technical Walkthrough: RTSP Ingestion ➔ Detection ➔ Tracking ➔ Recognition ➔ OSD Overlay"
    p_sub2.font.name = 'Calibri'
    p_sub2.font.size = Pt(14)
    p_sub2.font.color.rgb = MUTED_TEXT

    notes_slide1 = slide1.notes_slide
    notes_slide1.notes_text_frame.text = (
        "SPEAKER NOTES:\n"
        "Good morning/afternoon. Today I am going to walk you through our NVIDIA DeepStream hardware-accelerated surveillance pipeline. "
        "We have recently completed a total cleanup of our workspace, completely separating our Python GStreamer-bound pipeline "
        "and our C++ Native pipeline into 100% independent directories.\n\n"
        "This presentation will focus exclusively on how the frame data travels from the RTSP camera streams, "
        "through our neural network engines, matches identities against our database, and draws overlays using On-Screen Display (OSD)."
    )

    # ==========================================
    # SLIDE 2: High-Level Pipeline Data Flow
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_title(slide2, "DeepStream Pipeline: RTSP Ingestion to OSD Rendering")

    box_main2 = slide2.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.83), Inches(5.0))
    tf_main2 = box_main2.text_frame
    tf_main2.word_wrap = True

    steps = [
        ("1. Stream Ingestion (rtspsrc ➔ nvv4l2decoder)", "H.264 camera feeds are captured and decoded natively on the GPU's hardware decoder."),
        ("2. Stream Multiplexing (nvstreammux)", "Frames from multiple RTSP streams are batch-assembled into a single composite memory buffer."),
        ("3. Primary Inference: Face Detection (nvinfer - SCRFD)", "Hardware-accelerated face detection scans the batch and extracts raw coordinates and landmarks."),
        ("4. Bounding Box Decoding (libnvdsinfer_custom_impl_scrfd.so)", "A custom compiled C++ parser decodes the SCRFD output tensors into bounding box coordinates."),
        ("5. Object Tracking (nvtracker - NvDCF)", "NVIDIA's DCF tracker links face boxes across frames, assigning persistent 'Track IDs' to each person."),
        ("6. Secondary Inference: Embedding (nvinfer - AdaFace)", "Passes tracked face crops through AdaFace to extract 512-dimension vector representations."),
        ("7. Overlays & Rendering (nvdsosd)", "Labels recognized names/known statuses (with green boxes) or UNKNOWN tags (with red boxes) onto the video frame.")
    ]

    for title, desc in steps:
        p = tf_main2.add_paragraph()
        p.text = f"➔ {title}: "
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(13)
        run.font.color.rgb = LIGHT_TEXT
        p.space_after = Pt(6)

    notes_slide2 = slide2.notes_slide
    notes_slide2.notes_text_frame.text = (
        "SPEAKER NOTES:\n"
        "Let's trace the journey of a single frame. The camera streams are ingested via RTSP. "
        "The GStreamer RTSP source elements route the frames to 'nvv4l2decoder', which handles the video decoding directly on the GPU hardware. "
        "The stream multiplexer then bundles the decoded frames from all active camera sources into a single memory batch.\n\n"
        "Next, our Primary Inference element runs the SCRFD face detector. Because GStreamer needs to convert model tensors "
        "into bounding box coordinates, we compile a custom C++ shared library parser. Once the face boxes are decoded, the NvDCF Tracker "
        "assigns tracking IDs. These tracked face crops are passed through the secondary engine, AdaFace, which outputs 512-dimensional vectors. "
        "Finally, the On-Screen Display (OSD) overlay element renders these labels onto the display window."
    )

    # ==========================================
    # SLIDE 3: Isolated Directory Structure
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_title(slide3, "Isolated Folders: Actual File Structure Listings")

    box_left3 = slide3.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_left3 = box_left3.text_frame
    tf_left3.word_wrap = True
    
    p_l3 = tf_left3.paragraphs[0]
    p_l3.text = "python_deepstream_pipeline/"
    p_l3.font.name = 'Calibri'
    p_l3.font.size = Pt(22)
    p_l3.font.bold = True
    p_l3.font.color.rgb = ACCENT_GREEN
    p_l3.space_after = Pt(10)

    py_struct = (
        "├── main.py\n"
        "├── README.md\n"
        "├── configs/\n"
        "│   ├── config_infer_primary.txt\n"
        "│   ├── config_infer_secondary.txt\n"
        "│   ├── config_tracker_NvDCF_perf.yml\n"
        "│   ├── gallery_embeddings.txt\n"
        "│   ├── ha_meem_master_config.txt\n"
        "│   ├── labels_primary.txt\n"
        "│   └── labels_secondary.txt\n"
        "└── custom_parser/\n"
        "    ├── compile_flags.txt\n"
        "    ├── nvdsinfer_custom_parser_scrfd.cpp\n"
        "    ├── nvdsinfer_custom_parser_adaface.cpp\n"
        "    └── include/\n"
        "        ├── nvdsinfer_custom_impl.h\n"
        "        └── ..."
    )
    p_pys = tf_left3.add_paragraph()
    p_pys.text = py_struct
    p_pys.font.name = 'Courier New'
    p_pys.font.size = Pt(11)
    p_pys.font.color.rgb = LIGHT_TEXT

    box_right3 = slide3.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_right3 = box_right3.text_frame
    tf_right3.word_wrap = True
    
    p_r3 = tf_right3.paragraphs[0]
    p_r3.text = "cpp_deepstream_pipeline/"
    p_r3.font.name = 'Calibri'
    p_r3.font.size = Pt(22)
    p_r3.font.bold = True
    p_r3.font.color.rgb = ACCENT_BLUE
    p_r3.space_after = Pt(10)

    cpp_struct = (
        "├── README.md\n"
        "├── configs/\n"
        "│   ├── config_infer_primary.txt\n"
        "│   ├── config_infer_secondary.txt\n"
        "│   ├── config_tracker_NvDCF_perf.yml\n"
        "│   ├── gallery_embeddings.txt\n"
        "│   ├── ha_meem_master_config.txt\n"
        "│   ├── labels_primary.txt\n"
        "│   └── labels_secondary.txt\n"
        "└── custom_parser/\n"
        "    ├── compile_flags.txt\n"
        "    ├── nvdsinfer_custom_parser_scrfd.cpp\n"
        "    ├── nvdsinfer_custom_parser_adaface.cpp\n"
        "    └── include/\n"
        "        ├── nvdsinfer_custom_impl.h\n"
        "        └── ..."
    )
    p_cpps = tf_right3.add_paragraph()
    p_cpps.text = cpp_struct
    p_cpps.font.name = 'Courier New'
    p_cpps.font.size = Pt(11)
    p_cpps.font.color.rgb = LIGHT_TEXT

    notes_slide3 = slide3.notes_slide
    notes_slide3.notes_text_frame.text = (
        "SPEAKER NOTES:\n"
        "To prevent any code overlap or dependency issues, we have split the DeepStream implementation into two separate "
        "independent folders at the root level of our workspace.\n\n"
        "In 'python_deepstream_pipeline/', we have all the scripts, configs, and custom parser files required to run the GStreamer Python bindings. "
        "In 'cpp_deepstream_pipeline/', we have the exact native C++ configurations and parser sources. "
        "If we delete one folder, the other remains fully functional. There is zero directory coupling."
    )

    # ==========================================
    # SLIDE 4: C++ Native Pipeline (Option A)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_title(slide4, "C++ Native Pipeline: Configuration and Operation")

    box_main4 = slide4.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.5))
    tf_main4 = box_main4.text_frame
    tf_main4.word_wrap = True

    cpp_points = [
        ("Binary-Driven Execution", "The C++ pipeline is executed by feeding the config files directly into the compiled native deepstream-app binary. No Python interpreter is used."),
        ("Config Parsing (ha_meem_master_config.txt)", "Serves as the main blueprint, declaring sources, tiled display sizes, tracker files, and sub-infer properties."),
        ("C++ Custom Bounding Box Parser", "SCRFD outputs anchors across three strides. We compile nvdsinfer_custom_parser_scrfd.cpp into libnvdsinfer_custom_impl_scrfd.so to decode the spatial bounding boxes on-the-fly."),
        ("Compiled Embeddings Recognition", "The custom C++ parser (nvdsinfer_custom_parser_adaface.cpp) reads gallery_embeddings.txt into memory at runtime and maps vectors to identities natively inside C++.")
    ]

    for title, desc in cpp_points:
        p = tf_main4.add_paragraph()
        p.text = f"✔ {title}: "
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(14)
        run.font.color.rgb = LIGHT_TEXT
        p.space_after = Pt(10)

    notes_slide4 = slide4.notes_slide
    notes_slide4.notes_text_frame.text = (
        "SPEAKER NOTES:\n"
        "The C++ native pipeline (Option A) is driven entirely by configuration files. "
        "We execute it using the pre-compiled 'deepstream-app' binary. "
        "The master configuration maps all our inputs and references sub-configuration files for inference.\n\n"
        "To handle our face models, the GStreamer elements load our custom-built libraries. "
        "The SCRFD bounding box parser is compiled into a shared object. "
        "The AdaFace recognizer custom parser loads the gallery database from a flat text file 'gallery_embeddings.txt' "
        "and handles matching calculations inside C++ natively."
    )

    # ==========================================
    # SLIDE 5: Python GStreamer Pipeline (Option B)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_title(slide5, "Python Pipeline: Custom Control & FAISS Integration")

    box_main5 = slide5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.5))
    tf_main5 = box_main5.text_frame
    tf_main5.word_wrap = True

    py_points = [
        ("Python GStreamer Bindings", "Elements are initialized programmatically in Python using GObject bindings (e.g. Gst.ElementFactory.make) and assembled into a pipeline."),
        ("Metadata Pad Probe Callbacks", "Attaches a Python callback probe (sgie_src_pad_probe) to the source pad of the secondary GIE to capture tensor data as it flows through."),
        ("FAISS Database Integration", "Embeddings are extracted from the buffer metadata and matched against a binary FAISS database (gallery_embeddings_80px_faiss.npy) using cosine similarity."),
        ("Temporal Consensus Filters", "Maintains an embedding history window for each tracked face. This ensures that transient noise or lighting shifts do not trigger incorrect matching decisions.")
    ]

    for title, desc in py_points:
        p = tf_main5.add_paragraph()
        p.text = f"✔ {title}: "
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN
        
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(14)
        run.font.color.rgb = LIGHT_TEXT
        p.space_after = Pt(10)

    notes_slide5 = slide5.notes_slide
    notes_slide5.notes_text_frame.text = (
        "SPEAKER NOTES:\n"
        "The Python pipeline (Option B) constructs the same GStreamer pipeline but wraps it inside a Python process. "
        "This allows us to leverage python-level libraries like NumPy and FAISS.\n\n"
        "We attach a 'pad probe' callback on the secondary inference Gst.Pad. "
        "Every time a batch is processed, GStreamer invokes our Python function. "
        "This function extracts the generated face embeddings and matches them against our binary FAISS file. "
        "We also apply our temporal consensus aggregator, meaning we collect face signatures across multiple frames "
        "before making a final decision, drastically improving matching accuracy."
    )

    # ==========================================
    # SLIDE 6: Zero-Copy Memory & GPU Optimization
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_title(slide6, "Pipeline Optimization: GPU Zero-Copy Memory Flow")

    box_main6 = slide6.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.5))
    tf_main6 = box_main6.text_frame
    tf_main6.word_wrap = True

    opt_points = [
        ("NVIDIA Video Memory Management (NVMM)", "The multiplexer and inference elements are configured to use cuda-device buffers (nvbuf-memory-type=0). The video frames stay on the GPU memory during all pipeline operations."),
        ("No Host CPU Memory Copies", "Frame matrices are never copied into the system's host RAM (CPU memory) during ingestion, decoding, face detection, tracking, or OSD overlay rendering."),
        ("Metadata-Only Tensor Parsing", "The custom parsers only copy lightweight tensor data (like bounding boxes and 512-dimension float arrays) to system memory, avoiding large memory copies."),
        ("Lazy Frame Extraction", "If a recognized status changes or an alert is triggered, GStreamer performs a lazy-copy of that specific frame buffer to the CPU for saving snapshots, keeping the rest of the stream on the GPU.")
    ]

    for title, desc in opt_points:
        p = tf_main6.add_paragraph()
        p.text = f"✦ {title}: "
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(14)
        run.font.color.rgb = LIGHT_TEXT
        p.space_after = Pt(10)

    notes_slide6 = slide6.notes_slide
    notes_slide6.notes_text_frame.text = (
        "SPEAKER NOTES:\n"
        "A key detail of both pipelines is their performance optimization. "
        "By setting 'nvbuf-memory-type=0', we tell GStreamer to use NVIDIA Video Memory (NVMM).\n\n"
        "Because of this, the video frames stay entirely inside the GPU VRAM as they are decoded, scaled, "
        "and analyzed. The only things passed to the CPU are lightweight metadata—such as bounding box coordinates "
        "and the 512-dimension embedding vectors. The heavy video frame matrix is never copied to the CPU RAM, "
        "except when we lazily copy a single snapshot frame during an alert trigger. This ensures maximum throughput."
    )

    # ==========================================
    # SLIDE 7: Summary & Multi-Batch Configuration
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_title(slide7, "Summary of Engine Configurations")

    box_main7 = slide7.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.5))
    tf_main7 = box_main7.text_frame
    tf_main7.word_wrap = True

    configs = [
        ("Multi-Stream Multiplexing", "Configured to multiplex three cameras in sync with a live-source push timeout of 40ms, preventing lag."),
        ("Primary Detector Batch Size", "Face detector (SCRFD) is configured with batch-size=3. TensorRT processes three frames (one from each camera) concurrently in a single GPU inference pass."),
        ("Secondary Classifier Batch Size", "The AdaFace embedding model is configured with batch-size=16. This allows it to process up to 16 face chips from the three streams simultaneously in one step."),
        ("Static Model Input Profiles", "The SCRFD ONNX input was compiled to static 640x640 dimensions. This fixes the dynamic TensorRT engine building failures and ensures stable GPU compilation.")
    ]

    for title, desc in configs:
        p = tf_main7.add_paragraph()
        p.text = f"✔ {title}: "
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN
        
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(14)
        run.font.color.rgb = LIGHT_TEXT
        p.space_after = Pt(10)

    notes_slide7 = slide7.notes_slide
    notes_slide7.notes_text_frame.text = (
        "SPEAKER NOTES:\n"
        "To summarize our engine configurations, the multiplexer is tuned to receive three active camera streams. "
        "For efficiency, the primary detector (SCRFD) is batch-configured to process exactly 3 frames in a single pass. "
        "The secondary classifier (AdaFace) operates on a batch size of 16, allowing it to handle situations "
        "where multiple faces are detected across the streams simultaneously without bottlenecking the pipeline."
    )

    # ==========================================
    # SLIDE 8: Setup and Run Execution
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_title(slide8, "Independent Pipeline Compile & Run Commands")

    box_left8 = slide8.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_left8 = box_left8.text_frame
    tf_left8.word_wrap = True

    p_l8 = tf_left8.paragraphs[0]
    p_l8.text = "Python DeepStream Execution"
    p_l8.font.name = 'Calibri'
    p_l8.font.size = Pt(20)
    p_l8.font.bold = True
    p_l8.font.color.rgb = ACCENT_GREEN
    p_l8.space_after = Pt(8)

    py_commands = (
        "# 1. Compile Parsers in Python dir\n"
        "cd /app/python_deepstream_pipeline/custom_parser\n"
        "g++ -shared -fPIC -o libnvdsinfer_custom_impl_scrfd.so nvdsinfer_custom_parser_scrfd.cpp -Iinclude -I/opt/nvidia/deepstream/deepstream/sources/includes -I/usr/local/cuda/include -O3 -std=c++17\n"
        "g++ -shared -fPIC -o libnvdsinfer_custom_impl_adaface.so nvdsinfer_custom_parser_adaface.cpp -Iinclude -I/opt/nvidia/deepstream/deepstream/sources/includes -I/usr/local/cuda/include -O3 -std=c++17\n\n"
        "# 2. Start the Pipeline\n"
        "python3 /app/python_deepstream_pipeline/main.py"
    )
    p_pyc = tf_left8.add_paragraph()
    p_pyc.text = py_commands
    p_pyc.font.name = 'Courier New'
    p_pyc.font.size = Pt(10)
    p_pyc.font.color.rgb = LIGHT_TEXT

    box_right8 = slide8.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5))
    tf_right8 = box_right8.text_frame
    tf_right8.word_wrap = True

    p_r8 = tf_right8.paragraphs[0]
    p_r8.text = "C++ Native Execution"
    p_r8.font.name = 'Calibri'
    p_r8.font.size = Pt(20)
    p_r8.font.bold = True
    p_r8.font.color.rgb = ACCENT_BLUE
    p_r8.space_after = Pt(8)

    cpp_commands = (
        "# 1. Compile Parsers in C++ dir\n"
        "cd /app/cpp_deepstream_pipeline/custom_parser\n"
        "g++ -shared -fPIC -o libnvdsinfer_custom_impl_scrfd.so nvdsinfer_custom_parser_scrfd.cpp -Iinclude -I/opt/nvidia/deepstream/deepstream/sources/includes -I/usr/local/cuda/include -O3 -std=c++17\n"
        "g++ -shared -fPIC -o libnvdsinfer_custom_impl_adaface.so nvdsinfer_custom_parser_adaface.cpp -Iinclude -I/opt/nvidia/deepstream/deepstream/sources/includes -I/usr/local/cuda/include -O3 -std=c++17\n\n"
        "# 2. Start deepstream-app\n"
        "cd /app/cpp_deepstream_pipeline/configs\n"
        "deepstream-app -c ha_meem_master_config.txt"
    )
    p_cppc = tf_right8.add_paragraph()
    p_cppc.text = cpp_commands
    p_cppc.font.name = 'Courier New'
    p_cppc.font.size = Pt(10)
    p_cppc.font.color.rgb = LIGHT_TEXT

    notes_slide8 = slide8.notes_slide
    notes_slide8.notes_text_frame.text = (
        "SPEAKER NOTES:\n"
        "Finally, here are the step-by-step commands to run the pipelines inside our DeepStream Docker container. "
        "Since the directories are completely separated, you compile the C++ parsers within their respective directories. "
        "For Python, you compile the custom parsers in '/app/python_deepstream_pipeline/custom_parser/' and launch using 'python3'. "
        "For C++, you compile the custom parsers in '/app/cpp_deepstream_pipeline/custom_parser/' and launch the compiled binary "
        "using 'deepstream-app'. Both are fully independent and operational."
    )

    # ==========================================
    # SLIDE 9: Downstream Flow (API & Alerting Bot)
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_title(slide9, "Post-Pipeline: Downstream Event Flow & WhatsApp Alerts")

    box_main9 = slide9.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.83), Inches(4.5))
    tf_main9 = box_main9.text_frame
    tf_main9.word_wrap = True

    downstream_points = [
        ("Structured Local Outputs", "The active GStreamer pipeline records events to 'logs/events.jsonl' and crops face snapshots to the 'snapshots/' directory in real-time."),
        ("FastAPI Server (apps/api_server/main.py)", "Exposes the local events log over HTTP, offering REST endpoints and a Server-Sent Events (SSE) stream for real-time dashboard updates."),
        ("WhatsApp Bot (apps/alert_bot/whatsapp_bot.py)", "A background worker that polls the API server, compresses snapshot crops, and drives WhatsApp Web using Selenium to send instant alerts."),
        ("Asynchronous Post-Inference Lifecycle", "Decoupling the pipelines from the alert handler ensures that alerts are queued and dispatched safely without blocking or slowing down the primary video inference thread.")
    ]

    for title, desc in downstream_points:
        p = tf_main9.add_paragraph()
        p.text = f"✦ {title}: "
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN
        
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(14)
        run.font.color.rgb = LIGHT_TEXT
        p.space_after = Pt(10)

    notes_slide9 = slide9.notes_slide
    notes_slide9.notes_text_frame.text = (
        "SPEAKER NOTES:\n"
        "After the DeepStream pipeline finishes running, the downstream flow begins. "
        "The pipeline saves face snapshot files and logs events locally. "
        "Our FastAPI server acts as the central data hub, tailing the log file and exposing event history via API. "
        "Finally, the Selenium-based WhatsApp Bot polls this API and automates WhatsApp Web to send immediate notifications. "
        "Because these services run asynchronously, they do not impact the high-performance video processing pipeline."
    )

    # Save presentation
    output_path = "/Users/ajmunna/Desktop/Workspace/TDI WorkSpace/Ha-Meem Group/ha_meem_ai_surveillance/DeepStream_Surveillance_Presentation.pptx"
    prs.save(output_path)
    print(f"[SUCCESS] Presentation saved to {output_path}")

if __name__ == '__main__':
    create_presentation()
